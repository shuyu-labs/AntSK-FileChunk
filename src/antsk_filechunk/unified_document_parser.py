"""
统一文档解析器
整合enhanced_parser和document_parser copy的功能，提供统一的接口
支持PDF、Word、Excel、PowerPoint文档的解析，包括表格和图片的提取
"""

import json
import logging
import mimetypes
import os
import re
import uuid
from pathlib import Path
from typing import Dict, List, Tuple, Optional, Union
from dataclasses import dataclass
from urllib.error import HTTPError, URLError
from urllib.parse import urlparse
from urllib.request import Request, urlopen
import docx
from docx.shared import Inches
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
import fitz  # PyMuPDF
from PIL import Image
import io
import chardet

# 导入新增的解析库
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import openpyxl
    from openpyxl import load_workbook
except ImportError:
    openpyxl = None

try:
    import xlrd
except ImportError:
    xlrd = None

logger = logging.getLogger(__name__)

@dataclass
class DocumentContent:
    """统一的文档内容数据结构"""
    paragraphs: List[Dict]  # 段落列表，每个段落包含内容和元数据
    tables: List[Dict]      # 表格列表（包含markdown格式）
    images: List[Dict]      # 图片信息列表
    metadata: Dict          # 文档元数据
    structure: Dict         # 文档结构信息
    markdown_content: str   # 完整的markdown内容
    file_info: Dict         # 文件基本信息

class UnifiedDocumentParser:
    """统一文档解析器主类"""
    
    def __init__(self, image_base_url: str = "http://localhost:8000", 
                 image_save_dir: str = "./static/images"):
        """
        初始化统一文档解析器
        
        Args:
            image_base_url: 图片访问的基础URL
            image_save_dir: 图片保存目录
        """
        self.image_base_url = image_base_url
        self.image_save_dir = Path(image_save_dir)
        self.image_save_dir.mkdir(parents=True, exist_ok=True)
        self.remote_parse_url = os.getenv("FILEEXT_API_URL", "https://fileext.antsk.cn/api/parse")
        self.remote_parse_api_key = os.getenv("FILEEXT_API_KEY", "").strip()
        self.remote_parse_timeout = int(os.getenv("FILEEXT_API_TIMEOUT", "120"))
        self.remote_parse_enabled = os.getenv("FILEEXT_API_ENABLED", "true").strip().lower() not in {
            "0",
            "false",
            "off",
            "no",
        }
        
        # 支持的文件格式（仅跨平台格式）
        self.supported_formats = {'.pdf', '.docx', '.doc', '.txt', '.xlsx', '.xls', '.pptx'}
        self.remote_supported_formats = {'.pdf', '.docx', '.doc', '.xlsx', '.xls', '.pptx'}
        
        # 支持的图片格式
        self.supported_image_formats = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp'}
        
    def parse_file(self, file_path: Union[str, Path]) -> DocumentContent:
        """
        解析文档文件（统一接口）
        
        Args:
            file_path: 文件路径
            
        Returns:
            统一的文档内容对象
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"文件不存在: {file_path}")
        
        file_extension = file_path.suffix.lower()
        
        if file_extension not in self.supported_formats:
            raise ValueError(f"不支持的文件格式: {file_extension}。支持的格式：{', '.join(self.supported_formats)}")
        
        logger.info(f"开始解析文件: {file_path} (格式: {file_extension})")
        
        try:
            if self.remote_parse_enabled and file_extension in self.remote_supported_formats:
                try:
                    logger.info("尝试使用远程 Markdown 转换服务: %s", self.remote_parse_url)
                    return self._parse_via_remote_api(file_path, file_extension)
                except Exception as remote_error:
                    logger.warning("远程 Markdown 转换失败，回退到本地解析: %s", remote_error)
                    if file_extension == '.doc':
                        raise
            # 根据文件类型调用相应的解析方法
            if file_extension == '.pdf':
                return self._parse_pdf_unified(file_path)
            elif file_extension == '.docx':
                return self._parse_docx_unified(file_path)
            elif file_extension == '.doc':
                raise ValueError("当前环境未启用 .doc 本地解析，请配置远程 Markdown 转换服务")
            elif file_extension == '.txt':
                return self._parse_txt_unified(file_path)
            elif file_extension == '.xlsx':
                return self._parse_xlsx_unified(file_path)
            elif file_extension == '.xls':
                return self._parse_xls_unified(file_path)
            elif file_extension == '.pptx':
                return self._parse_pptx_unified(file_path)
            else:
                raise ValueError(f"未实现的文件格式处理: {file_extension}")
                
        except Exception as e:
            logger.error(f"文件解析失败: {e}")
            raise
    
    def _parse_via_remote_api(self, file_path: Path, file_extension: str) -> DocumentContent:
        """Call the remote AntSK parser and adapt the markdown response."""
        response_payload = self._call_remote_markdown_api(file_path)

        if not response_payload.get("success", False):
            raise RuntimeError(response_payload.get("error") or "远程文档解析返回失败")

        markdown = (response_payload.get("markdown") or "").strip()
        if not markdown:
            raise RuntimeError("远程文档解析成功，但未返回 markdown 内容")

        images = self._normalize_remote_images(response_payload.get("images") or [])
        document_content = self._build_document_content_from_markdown(
            markdown=markdown,
            file_path=file_path,
            file_extension=file_extension,
            remote_images=images,
        )
        document_content.metadata.update(
            {
                "parser_source": "remote_api",
                "remote_file_type": response_payload.get("file_type", file_extension.lstrip(".")),
                "remote_images_count": response_payload.get("images_count", len(images)),
            }
        )
        return document_content

    def _call_remote_markdown_api(self, file_path: Path) -> Dict:
        """Upload a document to the remote parsing service."""
        boundary = f"----AntSKBoundary{uuid.uuid4().hex}"
        content_type = mimetypes.guess_type(file_path.name)[0] or "application/octet-stream"

        with file_path.open("rb") as input_file:
            file_bytes = input_file.read()

        body = bytearray()
        body.extend(f"--{boundary}\r\n".encode("utf-8"))
        body.extend(
            (
                f'Content-Disposition: form-data; name="file"; filename="{file_path.name}"\r\n'
                f"Content-Type: {content_type}\r\n\r\n"
            ).encode("utf-8")
        )
        body.extend(file_bytes)
        body.extend(f"\r\n--{boundary}--\r\n".encode("utf-8"))

        headers = {
            "Accept": "application/json",
            "Content-Type": f"multipart/form-data; boundary={boundary}",
        }
        if self.remote_parse_api_key:
            headers["X-API-KEY"] = self.remote_parse_api_key

        request = Request(
            self.remote_parse_url,
            data=bytes(body),
            headers=headers,
            method="POST",
        )

        try:
            with urlopen(request, timeout=self.remote_parse_timeout) as response:
                charset = response.headers.get_content_charset() or "utf-8"
                payload = response.read().decode(charset, errors="replace")
        except HTTPError as exc:
            error_body = exc.read().decode("utf-8", errors="replace")
            raise RuntimeError(f"远程文档解析 HTTP {exc.code}: {error_body}") from exc
        except URLError as exc:
            raise RuntimeError(f"远程文档解析连接失败: {exc}") from exc

        try:
            return json.loads(payload)
        except json.JSONDecodeError as exc:
            raise RuntimeError("远程文档解析返回了无法解析的 JSON") from exc

    def _normalize_remote_images(self, remote_images: List[Union[Dict, str]]) -> List[Dict]:
        """Normalize image metadata returned by the remote service."""
        normalized_images = []

        for index, image_item in enumerate(remote_images):
            if isinstance(image_item, dict):
                image_url = (
                    image_item.get("url")
                    or image_item.get("image_url")
                    or image_item.get("src")
                    or image_item.get("path")
                    or ""
                )
                filename = image_item.get("filename") or self._infer_image_filename(image_url, index)
                normalized_images.append({**image_item, "url": image_url, "filename": filename})
            elif isinstance(image_item, str):
                normalized_images.append(
                    {
                        "url": image_item,
                        "filename": self._infer_image_filename(image_item, index),
                    }
                )

        return normalized_images

    def _build_document_content_from_markdown(
        self,
        markdown: str,
        file_path: Path,
        file_extension: str,
        remote_images: Optional[List[Dict]] = None,
    ) -> DocumentContent:
        """Convert remote markdown into the local DocumentContent shape."""
        markdown_content = self._clean_markdown(markdown.splitlines())
        paragraphs = self._extract_paragraphs_from_markdown(markdown_content)
        tables = self._extract_tables_from_markdown(markdown_content)
        images = self._merge_markdown_images(markdown_content, remote_images or [])

        paragraphs = self._postprocess_paragraphs(paragraphs)
        structure = self._build_document_structure(paragraphs, tables, images)

        file_info = {
            "name": file_path.name,
            "size": file_path.stat().st_size,
            "format": file_extension.lstrip("."),
            "source": "remote_api",
        }
        metadata = {
            "title": file_path.stem,
            "format": file_extension.lstrip(".").upper(),
            "parser": "fileext.antsk.cn",
        }

        return DocumentContent(
            paragraphs=paragraphs,
            tables=tables,
            images=images,
            metadata=metadata,
            structure=structure,
            markdown_content=markdown_content,
            file_info=file_info,
        )

    def _extract_paragraphs_from_markdown(self, markdown_content: str) -> List[Dict]:
        """Extract headings and paragraph blocks from markdown."""
        paragraphs = []
        block_index = 0

        for block in re.split(r"\n\s*\n", markdown_content):
            block = block.strip()
            if not block:
                continue

            if self._is_markdown_table(block):
                continue

            text_lines = []
            for raw_line in block.splitlines():
                line = raw_line.strip()
                if not line or re.match(r"^!\[[^\]]*\]\([^)]+\)$", line):
                    continue
                text_lines.append(line)

            if not text_lines:
                continue

            if len(text_lines) == 1:
                heading_match = re.match(r"^(#{1,6})\s+(.*)$", text_lines[0])
                if heading_match:
                    paragraphs.append(
                        {
                            "content": heading_match.group(2).strip(),
                            "index": len(paragraphs),
                            "type": "heading",
                            "level": len(heading_match.group(1)),
                            "block_index": block_index,
                        }
                    )
                    block_index += 1
                    continue

            paragraphs.append(
                {
                    "content": "\n".join(text_lines),
                    "index": len(paragraphs),
                    "type": "paragraph",
                    "block_index": block_index,
                }
            )
            block_index += 1

        return paragraphs

    def _extract_tables_from_markdown(self, markdown_content: str) -> List[Dict]:
        """Extract markdown tables from markdown blocks."""
        tables = []

        for block in re.split(r"\n\s*\n", markdown_content):
            block = block.strip()
            if not block or not self._is_markdown_table(block):
                continue

            table_data = self._parse_markdown_table(block)
            if not table_data:
                continue

            markdown_lines = [line.rstrip() for line in block.splitlines() if line.strip()]
            tables.append(
                {
                    "index": len(tables),
                    "data": table_data,
                    "rows": len(table_data),
                    "cols": len(table_data[0]) if table_data else 0,
                    "type": "table",
                    "markdown": markdown_lines,
                }
            )

        return tables

    def _merge_markdown_images(self, markdown_content: str, remote_images: List[Dict]) -> List[Dict]:
        """Merge markdown image references with remote image metadata."""
        image_pattern = re.compile(r"!\[([^\]]*)\]\(([^)]+)\)")
        merged_images = []
        seen_urls = set()

        for index, match in enumerate(image_pattern.finditer(markdown_content)):
            alt_text = match.group(1).strip()
            image_url = match.group(2).strip()
            seen_urls.add(image_url)
            merged_images.append(
                {
                    "url": image_url,
                    "filename": alt_text or self._infer_image_filename(image_url, index),
                }
            )

        for index, remote_image in enumerate(remote_images):
            image_url = remote_image.get("url", "").strip()
            if not image_url or image_url in seen_urls:
                continue

            merged_images.append(
                {
                    **remote_image,
                    "filename": remote_image.get("filename") or self._infer_image_filename(image_url, index),
                }
            )

        return merged_images

    def _infer_image_filename(self, image_url: str, index: int) -> str:
        """Infer a reasonable image filename from the remote URL."""
        parsed = urlparse(image_url)
        name = Path(parsed.path).name
        return name or f"image_{index + 1}.png"

    def _parse_pdf_unified(self, file_path: Path) -> DocumentContent:
        """统一的PDF解析方法"""
        try:
            doc = fitz.open(str(file_path))
            paragraphs = []
            tables = []
            images = []
            markdown_lines = []
            
            total_pages = len(doc)
            logger.info(f"PDF文档共{total_pages}页")
            
            # 获取文档元数据
            metadata = self._extract_pdf_metadata(doc)
            
            for page_num in range(total_pages):
                page = doc.load_page(page_num)
                
                # 提取文本内容（使用enhanced_parser的逻辑）
                text = page.get_text()
                if text.strip():
                    page_paragraphs, page_markdown = self._process_pdf_text_unified(text, page_num)
                    paragraphs.extend(page_paragraphs)
                    markdown_lines.extend(page_markdown)
                
                # 提取表格（简化版）
                page_tables = self._extract_pdf_tables_unified(page, page_num)
                tables.extend(page_tables)
                
                # 提取图片
                page_images = self._extract_pdf_images_unified(page, page_num)
                if page_images:
                    images.extend(page_images)
                    # 在markdown中添加图片
                    for img_info in page_images:
                        markdown_lines.append(f"![图片]({img_info['url']})")
                        markdown_lines.append("")
            
            doc.close()
            
            # 后处理段落
            paragraphs = self._postprocess_paragraphs(paragraphs)
            
            # 构建文档结构
            structure = self._build_document_structure(paragraphs, tables, images)
            
            # 生成完整的markdown内容
            markdown_content = self._clean_markdown(markdown_lines)
            
            # 文件信息
            file_info = {
                "name": file_path.name,
                "size": file_path.stat().st_size,
                "format": "pdf",
                "pages": total_pages
            }
            
            return DocumentContent(
                paragraphs=paragraphs,
                tables=tables,
                images=images,
                metadata=metadata,
                structure=structure,
                markdown_content=markdown_content,
                file_info=file_info
            )
            
        except Exception as e:
            logger.error(f"PDF解析失败: {e}")
            raise
    
    def _parse_docx_unified(self, file_path: Path) -> DocumentContent:
        """统一的DOCX解析方法"""
        try:
            doc = docx.Document(str(file_path))
            paragraphs = []
            tables = []
            images = []
            markdown_lines = []
            
            total_paragraphs = len(doc.paragraphs)
            total_tables = len(doc.tables)
            
            logger.info(f"DOCX文档包含{total_paragraphs}个段落和{total_tables}个表格")
            
            # 获取文档元数据
            metadata = self._extract_docx_metadata(doc)
            
            # 先提取所有图片信息，建立图片映射
            all_images = self._extract_docx_images_unified(doc)
            image_map = {}  # 用于存储图片关系ID到图片信息的映射
            
            # 建立图片关系映射
            try:
                for rel_id, rel in doc.part.rels.items():
                    if "image" in rel.target_ref:
                        # 查找对应的图片信息
                        for img_info in all_images:
                            if img_info.get('rel_id') == rel_id or img_info.get('path', '').endswith(rel.target_ref.split('/')[-1]):
                                image_map[rel_id] = img_info
                                break
            except Exception as e:
                logger.warning(f"建立图片映射时出错: {str(e)}")
            
            # 处理段落（使用enhanced_parser的逻辑）
            for i, paragraph in enumerate(doc.paragraphs):
                try:
                    para_data, para_markdown = self._process_docx_paragraph_unified(paragraph, i)
                    if para_data:
                        paragraphs.append(para_data)
                    if para_markdown:
                        markdown_lines.extend(para_markdown)
                    
                    # 检查段落中是否包含图片
                    para_images = self._extract_paragraph_images(paragraph, image_map)
                    if para_images:
                        # 将图片添加到当前位置
                        for img_info in para_images:
                            if img_info not in images:  # 避免重复添加
                                images.append(img_info)
                                markdown_lines.append(f"![{img_info.get('filename', '图片')}]({img_info['url']})")
                                markdown_lines.append("")
                        
                except Exception as e:
                    logger.warning(f"处理第{i+1}个段落时出错: {str(e)}")
                    continue
            
            # 处理表格（使用enhanced_parser的表格转markdown功能）
            for i, table in enumerate(doc.tables):
                try:
                    table_data = self._extract_docx_table_unified(table, i)
                    if table_data:
                        tables.append(table_data)
                        # 添加表格的markdown到内容中
                        markdown_lines.append(f"\n**表格 {i+1}:**\n")
                        markdown_lines.extend(table_data.get('markdown', []))
                        markdown_lines.append("")
                except Exception as e:
                    logger.warning(f"处理第{i+1}个表格时出错: {str(e)}")
                    markdown_lines.append(f"\n*[表格 {i+1} 解析失败]*\n")
            
            # 添加剩余未处理的图片（如果有的话）
            try:
                remaining_images = [img for img in all_images if img not in images]
                if remaining_images:
                    markdown_lines.append("\n## 其他图片\n")
                    for img_info in remaining_images:
                        images.append(img_info)
                        markdown_lines.append(f"![{img_info.get('filename', '图片')}]({img_info['url']})")
                        markdown_lines.append("")
            except Exception as e:
                logger.warning(f"处理剩余图片时出错: {str(e)}")
            
            # 后处理段落
            paragraphs = self._postprocess_paragraphs(paragraphs)
            
            # 构建文档结构
            structure = self._build_document_structure(paragraphs, tables, images)
            
            # 生成完整的markdown内容
            markdown_content = self._clean_markdown(markdown_lines)
            
            # 文件信息
            file_info = {
                "name": file_path.name,
                "size": file_path.stat().st_size,
                "format": "docx",
                "paragraphs": total_paragraphs,
                "tables": total_tables
            }
            
            return DocumentContent(
                paragraphs=paragraphs,
                tables=tables,
                images=images,
                metadata=metadata,
                structure=structure,
                markdown_content=markdown_content,
                file_info=file_info
            )
            
        except Exception as e:
            logger.error(f"DOCX解析失败: {e}")
            raise
    
    def _parse_txt_unified(self, file_path: Path) -> DocumentContent:
        """统一的TXT解析方法"""
        try:
            # 检测文件编码
            with open(file_path, 'rb') as f:
                raw_data = f.read()
                encoding = chardet.detect(raw_data)['encoding'] or 'utf-8'
            
            # 读取文本内容
            with open(file_path, 'r', encoding=encoding) as f:
                content = f.read()
            
            # 按空行分割段落
            para_texts = re.split(r'\n\s*\n', content)
            
            paragraphs = []
            tables = []
            markdown_lines = []
            
            for i, para_text in enumerate(para_texts):
                para_text = para_text.strip()
                if para_text and len(para_text) > 10:
                    # 检查是否为表格内容
                    table_data = self._detect_and_parse_table_from_text(para_text)
                    if table_data:
                        # 创建表格对象
                        table_obj = {
                            'index': len(tables),
                            'data': table_data,
                            'rows': len(table_data),
                            'cols': len(table_data[0]) if table_data else 0,
                            'type': 'table',
                            'markdown': self._convert_list_to_markdown_table(table_data)
                        }
                        tables.append(table_obj)
                        
                        # 添加表格的markdown到内容中
                        markdown_lines.append(f"\n**表格 {len(tables)}:**\n")
                        markdown_lines.extend(table_obj['markdown'])
                        markdown_lines.append("")
                    else:
                        # 普通段落
                        paragraphs.append({
                            'content': para_text,
                            'index': i,
                            'type': 'paragraph',
                            'line_start': 0,  # 简化处理
                            'line_end': 0
                        })
                        markdown_lines.append(para_text)
                        markdown_lines.append("")
            
            metadata = {
                'title': file_path.stem,
                'format': 'TXT',
                'encoding': encoding,
                'size': len(content)
            }
            
            structure = self._build_document_structure(paragraphs, tables, [])
            
            # 生成完整的markdown内容
            markdown_content = self._clean_markdown(markdown_lines)
            
            # 文件信息
            file_info = {
                "name": file_path.name,
                "size": file_path.stat().st_size,
                "format": "txt",
                "encoding": encoding
            }
            
            return DocumentContent(
                paragraphs=paragraphs,
                tables=tables,
                images=[],
                metadata=metadata,
                structure=structure,
                markdown_content=markdown_content,
                file_info=file_info
            )
            
        except Exception as e:
            logger.error(f"文本文件解析失败: {e}")
            raise
    
    def _parse_xlsx_unified(self, file_path: Path) -> DocumentContent:
        """统一的XLSX解析方法"""
        try:
            if openpyxl is None:
                raise ImportError("需要安装openpyxl库来解析xlsx文件")
            
            workbook = load_workbook(filename=str(file_path), data_only=True)
            paragraphs = []
            tables = []
            images = []
            markdown_lines = []
            
            logger.info(f"XLSX文件包含{len(workbook.sheetnames)}个工作表")
            
            # 提取图片
            try:
                images = self._extract_xlsx_images_unified(str(file_path))
                if images:
                    logger.info(f"从XLSX文件中提取了{len(images)}张图片")
            except Exception as e:
                logger.warning(f"XLSX图片提取失败: {str(e)}")
            
            # 处理每个工作表
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                
                # 获取有数据的范围
                if sheet.max_row == 1 and sheet.max_column == 1:
                    continue  # 跳过空工作表
                
                # 转换为表格格式
                table_data = []
                for row in sheet.iter_rows(min_row=1, max_row=sheet.max_row, 
                                         min_col=1, max_col=sheet.max_column):
                    if any(cell.value is not None and str(cell.value).strip() for cell in row):
                        # 处理每个单元格，正确处理日期类型
                        row_data = []
                        for cell in row:
                            cell_value = self._format_excel_cell_value(cell)
                            row_data.append(cell_value)
                        table_data.append(row_data)
                
                if table_data:
                    # 创建表格对象
                    table_obj = {
                        'index': len(tables),
                        'sheet_name': sheet_name,
                        'data': table_data,
                        'rows': len(table_data),
                        'cols': len(table_data[0]) if table_data else 0,
                        'type': 'table',
                        'markdown': self._convert_list_to_markdown_table(table_data)
                    }
                    tables.append(table_obj)
                    
                    # 添加到markdown内容
                    markdown_lines.append(f"\n**工作表: {sheet_name}**\n")
                    markdown_lines.extend(table_obj['markdown'])
                    markdown_lines.append("")
            
            # 如果有图片，添加到markdown内容中
            if images:
                markdown_lines.append("\n## 图片")
                for img_info in images:
                    markdown_lines.append(f"![{img_info['filename']}]({img_info['url']})")
            
            # 生成完整的markdown内容
            markdown_content = self._clean_markdown(markdown_lines)
            
            metadata = {
                'title': file_path.stem,
                'format': 'XLSX',
                'sheets': len(workbook.sheetnames),
                'sheet_names': workbook.sheetnames
            }
            
            structure = self._build_document_structure(paragraphs, tables, images)
            
            # 文件信息
            file_info = {
                "name": file_path.name,
                "size": file_path.stat().st_size,
                "format": "xlsx",
                "sheets": len(workbook.sheetnames)
            }
            
            return DocumentContent(
                paragraphs=paragraphs,
                tables=tables,
                images=images,
                metadata=metadata,
                structure=structure,
                markdown_content=markdown_content,
                file_info=file_info
            )
            
        except Exception as e:
            logger.error(f"XLSX解析失败: {e}")
            raise
    
    def _format_xls_cell_value(self, cell, workbook) -> str:
        """
        格式化XLS单元格的值，正确处理日期、时间等特殊类型
        
        Args:
            cell: xlrd的Cell对象
            workbook: xlrd的Workbook对象
            
        Returns:
            格式化后的字符串值
        """
        import datetime
        from xlrd import xldate_as_datetime, XL_CELL_DATE
        
        if cell.value is None or (isinstance(cell.value, str) and not cell.value.strip()):
            return ""
        
        # 检查是否为日期类型
        if cell.ctype == XL_CELL_DATE:
            try:
                # 将Excel日期序列号转换为datetime对象
                date_value = xldate_as_datetime(cell.value, workbook.datemode)
                
                # 判断是否只有日期部分
                if date_value.hour == 0 and date_value.minute == 0 and date_value.second == 0:
                    return date_value.strftime("%Y-%m-%d")
                else:
                    return date_value.strftime("%Y-%m-%d %H:%M:%S")
            except Exception as e:
                logger.warning(f"日期转换失败: {e}, 原始值: {cell.value}")
                return str(cell.value)
        
        # 数字类型
        if isinstance(cell.value, float):
            # 如果是整数，显示为整数
            if cell.value.is_integer():
                return str(int(cell.value))
            else:
                return str(cell.value)
        
        # 其他类型直接转换为字符串
        return str(cell.value).strip()
    
    def _format_excel_cell_value(self, cell) -> str:
        """
        格式化Excel单元格的值，正确处理日期、时间等特殊类型
        
        Args:
            cell: openpyxl的Cell对象
            
        Returns:
            格式化后的字符串值
        """
        import datetime
        
        if cell.value is None:
            return ""
        
        # 检查是否为日期或时间类型
        if isinstance(cell.value, (datetime.datetime, datetime.date, datetime.time)):
            if isinstance(cell.value, datetime.datetime):
                # 日期时间类型 - 如果时间部分为00:00:00，只显示日期
                if cell.value.hour == 0 and cell.value.minute == 0 and cell.value.second == 0:
                    return cell.value.strftime("%Y-%m-%d")
                else:
                    return cell.value.strftime("%Y-%m-%d %H:%M:%S")
            elif isinstance(cell.value, datetime.date):
                # 日期类型
                return cell.value.strftime("%Y-%m-%d")
            elif isinstance(cell.value, datetime.time):
                # 时间类型
                return cell.value.strftime("%H:%M:%S")
        
        # 数字类型 - 检查是否应该是日期（通过number_format判断）
        if isinstance(cell.value, (int, float)):
            # 检查单元格的数字格式
            if cell.number_format:
                # 常见的日期格式代码
                date_formats = ['yyyy', 'yy', 'mm', 'dd', 'm/d', 'd/m', 'h:mm', 'h:mm:ss']
                number_format_lower = cell.number_format.lower()
                
                # 如果格式中包含日期相关的标识
                if any(fmt in number_format_lower for fmt in date_formats):
                    try:
                        # 尝试将数字转换为日期
                        from openpyxl.utils.datetime import from_excel
                        date_value = from_excel(cell.value)
                        if isinstance(date_value, datetime.datetime):
                            return date_value.strftime("%Y-%m-%d %H:%M:%S")
                        elif isinstance(date_value, datetime.date):
                            return date_value.strftime("%Y-%m-%d")
                    except Exception as e:
                        logger.debug(f"日期转换失败: {e}")
            
            # 如果是整数且小数部分为0，显示为整数
            if isinstance(cell.value, float) and cell.value.is_integer():
                return str(int(cell.value))
        
        # 其他类型直接转换为字符串
        return str(cell.value).strip()
    
    def _parse_xls_unified(self, file_path: Path) -> DocumentContent:
        """统一的XLS解析方法"""
        try:
            if xlrd is None:
                raise ImportError("需要安装xlrd库来解析xls文件")
            
            # 尝试打开带格式信息的工作簿，如果失败则不带格式信息打开
            try:
                workbook = xlrd.open_workbook(str(file_path), formatting_info=True)
            except Exception:
                logger.warning("无法获取格式信息，使用基础模式打开XLS文件")
                workbook = xlrd.open_workbook(str(file_path))
            
            paragraphs = []
            tables = []
            images = []
            markdown_lines = []
            
            logger.info(f"XLS文件包含{workbook.nsheets}个工作表")
            
            # 处理每个工作表
            for sheet_idx in range(workbook.nsheets):
                sheet = workbook.sheet_by_index(sheet_idx)
                sheet_name = workbook.sheet_names()[sheet_idx]
                
                if sheet.nrows == 0:
                    continue  # 跳过空工作表
                
                # 转换为表格格式
                table_data = []
                for row_idx in range(sheet.nrows):
                    row_data = []
                    for col_idx in range(sheet.ncols):
                        cell = sheet.cell(row_idx, col_idx)
                        cell_value = self._format_xls_cell_value(cell, workbook)
                        row_data.append(cell_value)
                    
                    if any(cell.strip() for cell in row_data):
                        table_data.append(row_data)
                
                if table_data:
                    # 创建表格对象
                    table_obj = {
                        'index': len(tables),
                        'sheet_name': sheet_name,
                        'data': table_data,
                        'rows': len(table_data),
                        'cols': len(table_data[0]) if table_data else 0,
                        'type': 'table',
                        'markdown': self._convert_list_to_markdown_table(table_data)
                    }
                    tables.append(table_obj)
                    
                    # 添加到markdown内容
                    markdown_lines.append(f"\n**工作表: {sheet_name}**\n")
                    markdown_lines.extend(table_obj['markdown'])
                    markdown_lines.append("")
            
            # 生成完整的markdown内容
            markdown_content = self._clean_markdown(markdown_lines)
            
            metadata = {
                'title': file_path.stem,
                'format': 'XLS',
                'sheets': workbook.nsheets,
                'sheet_names': workbook.sheet_names()
            }
            
            structure = self._build_document_structure(paragraphs, tables, images)
            
            # 文件信息
            file_info = {
                "name": file_path.name,
                "size": file_path.stat().st_size,
                "format": "xls",
                "sheets": workbook.nsheets
            }
            
            return DocumentContent(
                paragraphs=paragraphs,
                tables=tables,
                images=images,
                metadata=metadata,
                structure=structure,
                markdown_content=markdown_content,
                file_info=file_info
            )
            
        except Exception as e:
            logger.error(f"XLS解析失败: {e}")
            raise
    
    def _parse_pptx_unified(self, file_path: Path) -> DocumentContent:
        """统一的PPTX解析方法"""
        try:
            if Presentation is None:
                raise ImportError("需要安装python-pptx库来解析pptx文件")
            
            prs = Presentation(str(file_path))
            paragraphs = []
            tables = []
            images = []
            markdown_lines = []
            
            logger.info(f"PPTX文件包含{len(prs.slides)}张幻灯片")
            
            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_paragraphs = []
                slide_text = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content = shape.text.strip()
                        slide_text.append(text_content)
                        
                        # 创建段落对象
                        para_obj = {
                            'content': text_content,
                            'index': len(paragraphs),
                            'type': 'paragraph',
                            'slide': slide_idx,
                            'shape_type': 'text'
                        }
                        paragraphs.append(para_obj)
                        slide_paragraphs.append(para_obj)
                    
                    # 提取图片
                    if shape.shape_type == 13:  # PICTURE = 13
                        try:
                            image = shape.image
                            image_bytes = image.blob
                            image_ext = self._get_image_extension(image_bytes)
                            image_filename = f"{uuid.uuid4()}.{image_ext}"
                            image_path = self.image_save_dir / image_filename
                            
                            with open(image_path, 'wb') as f:
                                f.write(image_bytes)
                            
                            image_url = f"{self.image_base_url}/static/images/{image_filename}"
                            image_info = {
                                "filename": image_filename,
                                "path": str(image_path),
                                "url": image_url,
                                "size": len(image_bytes),
                                "slide": slide_idx,
                                "type": "image"
                            }
                            images.append(image_info)
                            
                            slide_text.append(f"![图片]({image_url})")
                            
                        except Exception as e:
                            logger.warning(f"提取PPT图片失败: {str(e)}")
                
                # 添加幻灯片内容到markdown
                if slide_text:
                    markdown_lines.append(f"\n## 幻灯片 {slide_idx}\n")
                    markdown_lines.extend(slide_text)
                    markdown_lines.append("")
            
            # 生成完整的markdown内容
            markdown_content = self._clean_markdown(markdown_lines)
            
            metadata = {
                'title': file_path.stem,
                'format': 'PPTX',
                'slides': len(prs.slides)
            }
            
            structure = self._build_document_structure(paragraphs, tables, images)
            
            # 文件信息
            file_info = {
                "name": file_path.name,
                "size": file_path.stat().st_size,
                "format": "pptx",
                "slides": len(prs.slides)
            }
            
            return DocumentContent(
                paragraphs=paragraphs,
                tables=tables,
                images=images,
                metadata=metadata,
                structure=structure,
                markdown_content=markdown_content,
                file_info=file_info
            )
            
        except Exception as e:
            logger.error(f"PPTX解析失败: {e}")
            raise
    
    # 以下是辅助方法，整合了两个解析器的功能
    
    def _process_pdf_text_unified(self, text: str, page_num: int) -> Tuple[List[Dict], List[str]]:
        """处理PDF文本，返回段落对象和markdown行"""
        paragraphs = []
        markdown_lines = []
        
        # 按段落分割
        para_texts = text.split('\n\n')
        
        for i, para in enumerate(para_texts):
            para = para.strip()
            if not para:
                continue
            
            # 清理多余的换行符
            para = re.sub(r'\n+', ' ', para)
            para = re.sub(r'\s+', ' ', para)
            
            if len(para) < 10:  # 过滤过短的段落
                continue
            
            # 检测标题
            is_heading = self._is_heading_pdf(para)
            
            # 创建段落对象
            para_obj = {
                'content': para,
                'page': page_num + 1,
                'index': i,
                'type': 'heading' if is_heading else 'paragraph',
                'bbox': []
            }
            paragraphs.append(para_obj)
            
            # 添加到markdown
            if is_heading:
                markdown_lines.append(f"## {para}")
            else:
                markdown_lines.append(para)
            markdown_lines.append("")
        
        return paragraphs, markdown_lines
    
    def _process_docx_paragraph_unified(self, paragraph, index: int) -> Tuple[Optional[Dict], List[str]]:
        """处理DOCX段落，返回段落对象和markdown行"""
        text = paragraph.text.strip()
        
        if not text:
            return None, [""]
        
        # 检查段落样式
        style_name = paragraph.style.name if paragraph.style else ""
        
        # 创建段落对象
        para_obj = {
            'content': text,
            'index': index,
            'type': 'paragraph',
            'style': style_name,
            'level': self._get_heading_level(style_name)
        }
        
        markdown_lines = []
        
        # 处理标题
        if style_name.startswith('Heading'):
            level = self._get_heading_level(style_name)
            para_obj['type'] = 'heading'
            markdown_lines.append(f"{'#' * level} {text}")
            markdown_lines.append("")
        elif style_name.lower() in ['title', 'subtitle']:
            para_obj['type'] = 'title'
            markdown_lines.append(f"# {text}")
            markdown_lines.append("")
        else:
            # 检查段落对齐方式
            alignment = paragraph.alignment
            if alignment == WD_PARAGRAPH_ALIGNMENT.CENTER:
                markdown_lines.append(f"<center>{text}</center>")
            else:
                # 处理列表项
                if re.match(r'^\s*[\d\w]+[.)]\s+', text) or text.strip().startswith(('•', '-', '*')):
                    if re.match(r'^\s*\d+[.)]\s+', text):
                        match = re.match(r'^\s*\d+[.)]\s+(.+)', text)
                        if match:
                            markdown_lines.append(f"1. {match.group(1)}")
                    else:
                        cleaned_text = re.sub(r'^\s*[•\-*]\s+', '', text)
                        markdown_lines.append(f"- {cleaned_text}")
                else:
                    markdown_lines.append(text)
            
            markdown_lines.append("")
        
        return para_obj, markdown_lines
    
    def _extract_docx_table_unified(self, table, index: int) -> Dict:
        """提取Word表格信息并转换为markdown"""
        table_data = []
        
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                row_data.append(cell_text)
            table_data.append(row_data)
        
        # 转换为markdown
        markdown_table = self._convert_list_to_markdown_table(table_data)
        
        return {
            'index': index,
            'data': table_data,
            'rows': len(table.rows),
            'cols': len(table.columns) if table.rows else 0,
            'type': 'table',
            'markdown': markdown_table
        }
    
    def _extract_pdf_tables_unified(self, page, page_num: int) -> List[Dict]:
        """从PDF页面提取表格（简化版）"""
        # 这里可以集成更专业的表格提取库
        return []
    
    def _extract_pdf_images_unified(self, page, page_num: int) -> List[Dict]:
        """从PDF页面提取图片"""
        images_info = []
        
        try:
            image_list = page.get_images()
            
            for img_index, img in enumerate(image_list):
                try:
                    xref = img[0]
                    pix = fitz.Pixmap(page.parent, xref)
                    
                    # 跳过CMYK图片和太小的图片
                    if pix.n - pix.alpha >= 4 or pix.width < 50 or pix.height < 50:
                        pix = None
                        continue
                    
                    # 生成唯一文件名
                    image_id = str(uuid.uuid4())
                    image_filename = f"{image_id}.png"
                    image_path = self.image_save_dir / image_filename
                    
                    # 保存图片
                    pix.save(str(image_path))
                    
                    # 验证图片是否有效
                    if self._validate_image(image_path):
                        image_url = f"{self.image_base_url}/static/images/{image_filename}"
                        
                        images_info.append({
                            "filename": image_filename,
                            "path": str(image_path),
                            "url": image_url,
                            "page": page_num + 1,
                            "index": img_index,
                            "width": pix.width,
                            "height": pix.height,
                            "format": "png",
                            "type": "image"
                        })
                    else:
                        os.remove(image_path)
                    
                    pix = None
                    
                except Exception as e:
                    logger.warning(f"处理PDF第{page_num + 1}页图片{img_index}时出错: {str(e)}")
                    continue
                    
        except Exception as e:
            logger.warning(f"PDF第{page_num + 1}页图片提取失败: {str(e)}")
        
        return images_info
    
    def _extract_docx_images_unified(self, doc) -> List[Dict]:
        """从docx文档中提取图片"""
        images_info = []
        
        try:
            for rel_id, rel in doc.part.rels.items():
                if "image" in rel.target_ref:
                    try:
                        image_data = rel.target_part.blob
                        
                        if len(image_data) < 100:
                            continue
                        
                        image_id = str(uuid.uuid4())
                        image_ext = self._get_image_extension(image_data)
                        image_filename = f"{image_id}.{image_ext}"
                        image_path = self.image_save_dir / image_filename
                        
                        with open(image_path, 'wb') as f:
                            f.write(image_data)
                        
                        if self._validate_image(image_path):
                            image_url = f"{self.image_base_url}/static/images/{image_filename}"
                            
                            images_info.append({
                                "filename": image_filename,
                                "path": str(image_path),
                                "url": image_url,
                                "size": len(image_data),
                                "format": image_ext,
                                "type": "image",
                                "rel_id": rel_id  # 添加关系ID用于映射
                            })
                        else:
                            os.remove(image_path)
                            
                    except Exception as e:
                        logger.warning(f"处理DOCX图片{rel_id}时出错: {str(e)}")
                        continue
                        
        except Exception as e:
            logger.warning(f"DOCX图片提取失败: {str(e)}")
        
        return images_info
    
    def _extract_paragraph_images(self, paragraph, image_map: Dict) -> List[Dict]:
        """从段落中提取图片信息"""
        para_images = []
        
        try:
            # 检查段落中的运行（runs）是否包含图片
            for run in paragraph.runs:
                # 检查运行中的绘图元素
                for drawing in run._element.xpath('.//a:blip'):
                    # 获取图片的关系ID
                    embed_attr = drawing.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                    if embed_attr and embed_attr in image_map:
                        para_images.append(image_map[embed_attr])
                
                # 也检查内联形状
                for inline in run._element.xpath('.//wp:inline'):
                    blips = inline.xpath('.//a:blip')
                    for blip in blips:
                        embed_attr = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                        if embed_attr and embed_attr in image_map:
                            para_images.append(image_map[embed_attr])
                            
        except Exception as e:
            logger.warning(f"提取段落图片时出错: {str(e)}")
        
        return para_images
    
    def _extract_xlsx_images_unified(self, file_path: str) -> List[Dict]:
        """从xlsx文件中提取图片"""
        images_info = []
        
        try:
            import zipfile
            
            with zipfile.ZipFile(file_path, 'r') as zip_file:
                media_files = [name for name in zip_file.namelist() 
                              if name.startswith('xl/media/') and 
                              any(name.lower().endswith(ext) for ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp'])]
                
                for media_file in media_files:
                    try:
                        image_data = zip_file.read(media_file)
                        
                        if len(image_data) < 100:
                            continue
                        
                        image_id = str(uuid.uuid4())
                        original_ext = Path(media_file).suffix.lower().lstrip('.')
                        image_filename = f"{image_id}.{original_ext}"
                        image_path = self.image_save_dir / image_filename
                        
                        with open(image_path, 'wb') as f:
                            f.write(image_data)
                        
                        if self._validate_image(image_path):
                            image_url = f"{self.image_base_url}/static/images/{image_filename}"
                            
                            images_info.append({
                                "filename": image_filename,
                                "path": str(image_path),
                                "url": image_url,
                                "size": len(image_data),
                                "format": original_ext,
                                "source": media_file,
                                "type": "image"
                            })
                        else:
                            os.remove(image_path)
                            
                    except Exception as e:
                        logger.warning(f"处理XLSX图片{media_file}时出错: {str(e)}")
                        
        except Exception as e:
            logger.error(f"XLSX图片提取失败: {str(e)}")
            
        return images_info
    
    def _convert_list_to_markdown_table(self, table_data: List[List[str]]) -> List[str]:
        """将二维列表转换为markdown表格"""
        if not table_data:
            return ["*空表格*"]
        
        markdown_table = []
        
        try:
            # 确保所有行有相同的列数，默认3列
            target_cols = 3
            normalized_data = []
            
            for row in table_data:
                # 清理和规范化每个单元格
                cleaned_row = []
                for i in range(target_cols):
                    if i < len(row):
                        cell = str(row[i]).strip()
                        # 清理单元格内容
                        cell = cell.replace('\n', ' ').replace('\r', ' ')
                        cell = re.sub(r'\s+', ' ', cell)  # 合并多个空格
                        # 转义markdown特殊字符
                        cell = cell.replace('|', '\\|')
                        cleaned_row.append(cell)
                    else:
                        cleaned_row.append('')
                normalized_data.append(cleaned_row)
            
            # 生成markdown表格
            if normalized_data:
                # 表头（第一行）
                headers = normalized_data[0]
                markdown_table.append("| " + " | ".join(headers) + " |")
                
                # 分隔线
                markdown_table.append("| " + " | ".join(["---"] * len(headers)) + " |")
                
                # 数据行（从第二行开始）
                for row in normalized_data[1:]:
                    markdown_table.append("| " + " | ".join(row) + " |")
            
            # 如果只有一行数据，添加默认表头
            if len(normalized_data) == 1:
                # 在前面添加默认表头
                default_headers = ['编号', '检查项', '描述']
                markdown_table.insert(0, "| " + " | ".join(["---"] * len(default_headers)) + " |")
                markdown_table.insert(0, "| " + " | ".join(default_headers) + " |")
            
        except Exception as e:
            logger.warning(f"表格转换失败: {str(e)}")
            return ["*表格解析失败*"]
        
        return markdown_table
    
    # 辅助方法
    def _validate_image(self, image_path: Path) -> bool:
        """验证图片是否有效"""
        try:
            with Image.open(image_path) as img:
                img.verify()
            return True
        except Exception:
            return False
    
    def _get_image_extension(self, image_data: bytes) -> str:
        """根据图片数据判断图片格式"""
        try:
            image = Image.open(io.BytesIO(image_data))
            format_name = image.format.lower()
            if format_name == 'jpeg':
                return 'jpg'
            return format_name if format_name else 'png'
        except Exception:
            return 'png'
    
    def _is_heading_pdf(self, text: str) -> bool:
        """判断PDF文本是否为标题"""
        if len(text) > 100:
            return False
        
        # 检测标题模式
        if (text.isupper() or 
            re.match(r'^[A-Z][^.!?]*$', text) or
            re.match(r'^\d+[\.\s]+[A-Z]', text)):
            return True
        
        return False
    
    def _get_heading_level(self, style_name: str) -> int:
        """获取标题级别"""
        heading_levels = {
            'Heading 1': 1, 'Title': 1,
            'Heading 2': 2, 'Subtitle': 2,
            'Heading 3': 3,
            'Heading 4': 4,
            'Heading 5': 5,
            'Heading 6': 6
        }
        return heading_levels.get(style_name, 0)
    
    def _clean_markdown(self, markdown_lines: List[str]) -> str:
        """清理和格式化markdown内容"""
        content = "\n".join(markdown_lines)
        content = re.sub(r'\n{3,}', '\n\n', content)
        lines = [line.rstrip() for line in content.split('\n')]
        
        # 移除开头和结尾的空行
        while lines and not lines[0].strip():
            lines.pop(0)
        while lines and not lines[-1].strip():
            lines.pop()
        
        return '\n'.join(lines)
    
    def _extract_pdf_metadata(self, doc) -> Dict:
        """提取PDF元数据"""
        metadata = doc.metadata or {}
        return {
            'title': metadata.get('title', ''),
            'author': metadata.get('author', ''),
            'subject': metadata.get('subject', ''),
            'creator': metadata.get('creator', ''),
            'producer': metadata.get('producer', ''),
            'creation_date': metadata.get('creationDate', ''),
            'modification_date': metadata.get('modDate', ''),
            'page_count': len(doc),
            'format': 'PDF'
        }
    
    def _extract_docx_metadata(self, doc) -> Dict:
        """提取Word文档元数据"""
        core_props = doc.core_properties
        return {
            'title': core_props.title or '',
            'author': core_props.author or '',
            'subject': core_props.subject or '',
            'keywords': core_props.keywords or '',
            'comments': core_props.comments or '',
            'created': str(core_props.created) if core_props.created else '',
            'modified': str(core_props.modified) if core_props.modified else '',
            'format': 'DOCX'
        }
    
    def _postprocess_paragraphs(self, paragraphs: List[Dict]) -> List[Dict]:
        """后处理段落列表"""
        processed = []
        
        for para in paragraphs:
            content = para['content']
            
            # 清理内容
            content = self._clean_paragraph_content(content)
            
            # 标题即使较短也需要保留，避免破坏文档结构
            if para.get('type') == 'heading' and content.strip():
                para['content'] = content
                processed.append(para)
                continue

            # 过滤空内容和过短的段落
            if not content.strip() or len(content.strip()) < 10:
                continue
            
            para['content'] = content
            processed.append(para)
        
        return processed
    
    def _clean_paragraph_content(self, content: str) -> str:
        """清理段落内容"""
        # 移除多余的空白字符
        content = re.sub(r'\s+', ' ', content)
        
        # 移除页眉页脚常见模式
        patterns_to_remove = [
            r'^第\s*\d+\s*页.*$',
            r'^\d+\s*$',
            r'^[-=_]{3,}$',
        ]
        
        for pattern in patterns_to_remove:
            if re.match(pattern, content.strip()):
                return ''
        
        return content.strip()
    
    def _build_document_structure(
        self, 
        paragraphs: List[Dict], 
        tables: List[Dict], 
        images: List[Dict]
    ) -> Dict:
        """构建文档结构信息"""
        structure = {
            'paragraph_count': len(paragraphs),
            'table_count': len(tables),
            'image_count': len(images),
            'headings': [],
            'sections': []
        }
        
        # 提取标题信息
        current_section = {'title': '', 'level': 0, 'start_index': 0, 'paragraphs': []}
        
        for i, para in enumerate(paragraphs):
            if para.get('type') == 'heading':
                if current_section['paragraphs']:
                    structure['sections'].append(current_section)
                
                level = para.get('level', 1)
                current_section = {
                    'title': para['content'],
                    'level': level,
                    'start_index': i,
                    'paragraphs': []
                }
                
                structure['headings'].append({
                    'text': para['content'],
                    'level': level,
                    'index': i
                })
            else:
                current_section['paragraphs'].append(i)
        
        # 保存最后一个章节
        if current_section['paragraphs']:
            structure['sections'].append(current_section)
        
        return structure
    
    def _detect_and_parse_table_from_text(self, text: str) -> Optional[List[List[str]]]:
        """
        从文本中检测和解析表格内容
        
        支持的格式：
        1. [表格内容] ... | ... | ...
        2. 用 | 分隔的表格行
        3. 标准markdown表格
        
        Args:
            text: 输入文本
            
        Returns:
            表格数据（二维列表）或None
        """
        try:
            # 格式1: [表格内容] ... | ... | ...
            if text.startswith('[表格内容]'):
                # 移除[表格内容]标记
                table_text = text.replace('[表格内容]', '').strip()
                return self._parse_pipe_separated_table(table_text)
            
            # 格式2: 检查是否包含多个 | 分隔符（可能是表格）
            if text.count('|') >= 3:  # 至少3个分隔符才可能是表格
                return self._parse_pipe_separated_table(text)
            
            # 格式3: 标准markdown表格
            if self._is_markdown_table(text):
                return self._parse_markdown_table(text)
            
            return None
            
        except Exception as e:
            logger.warning(f"表格解析失败: {e}")
            return None
    
    def _parse_pipe_separated_table(self, text: str) -> Optional[List[List[str]]]:
        """解析用管道符分隔的表格"""
        try:
            # 按 | 分割，保留空字符串以维持表格结构
            parts = [part.strip() for part in text.split('|')]
            
            # 移除首尾的空字符串
            while parts and not parts[0]:
                parts.pop(0)
            while parts and not parts[-1]:
                parts.pop()
            
            if len(parts) < 3:  # 至少需要3个部分才能构成表格
                return None
            
            table_data = []
            
            # 检查第一部分是否包含表头信息
            first_part = parts[0]
            has_header = any(keyword in first_part for keyword in ['编号', '检查项', '描述', '项目', '内容', '序号'])
            
            if has_header:
                # 提取表头
                header_parts = [p.strip() for p in first_part.split() if p.strip()]
                if len(header_parts) >= 2:
                    # 确保表头有3列
                    if len(header_parts) == 2:
                        header_parts.append('描述')
                    table_data.append(header_parts[:3])
                    parts = parts[1:]  # 移除已处理的表头部分
            
            # 如果没有表头，添加默认表头
            if not table_data:
                table_data.append(['编号', '检查项', '描述'])
            
            # 智能解析数据行
            current_row = []
            i = 0
            
            while i < len(parts):
                part = parts[i].strip()
                
                # 检查是否为编号开头
                number_match = re.match(r'^(\d+)\s*(.*)$', part)
                if number_match:
                    # 保存之前的行
                    if current_row and len(current_row) >= 2:
                        # 确保行有3列
                        while len(current_row) < 3:
                            current_row.append('')
                        table_data.append(current_row[:3])
                    
                    # 开始新行
                    number = number_match.group(1)
                    remaining_text = number_match.group(2).strip()
                    
                    # 尝试分离问题和描述
                    if remaining_text:
                        # 查找问号位置来分离问题和描述
                        question_end = remaining_text.find('?')
                        if question_end > 0:
                            question = remaining_text[:question_end + 1].strip()
                            description = remaining_text[question_end + 1:].strip()
                            # 移除可能的"示例:"前缀
                            description = re.sub(r'^\s*示例[：:]\s*', '', description)
                            current_row = [number, question, description]
                        else:
                            # 没有问号，尝试按空格分割
                            parts_split = remaining_text.split(None, 1)
                            if len(parts_split) >= 2:
                                current_row = [number, parts_split[0], parts_split[1]]
                            else:
                                current_row = [number, remaining_text, '']
                    else:
                        # 编号后面没有内容，等待下一个部分
                        current_row = [number]
                
                elif current_row:
                    # 继续填充当前行
                    if len(current_row) == 1:  # 只有编号
                        current_row.append(part)
                    elif len(current_row) == 2:  # 有编号和问题
                        current_row.append(part)
                    else:  # 已经有3列，合并到描述中
                        current_row[2] += ' ' + part
                else:
                    # 没有当前行，可能是独立的内容
                    if part:  # 非空内容
                        current_row = [str(len(table_data)), part, '']
                
                i += 1
            
            # 保存最后一行
            if current_row and len(current_row) >= 2:
                while len(current_row) < 3:
                    current_row.append('')
                table_data.append(current_row[:3])
            
            # 如果解析结果不理想，尝试简单的三列分割
            if len(table_data) < 2:
                table_data = [['编号', '检查项', '描述']]
                
                # 按3个一组分割原始parts
                for i in range(0, len(parts), 3):
                    row = parts[i:i+3]
                    if len(row) >= 1 and any(cell.strip() for cell in row):
                        # 补齐到3列
                        while len(row) < 3:
                            row.append('')
                        # 清理每个单元格
                        cleaned_row = [cell.strip() for cell in row[:3]]
                        table_data.append(cleaned_row)
            
            # 最终验证和清理
            if len(table_data) >= 2:
                # 确保所有行都有相同的列数
                max_cols = 3
                cleaned_table = []
                for row in table_data:
                    cleaned_row = [str(cell).strip() for cell in row[:max_cols]]
                    while len(cleaned_row) < max_cols:
                        cleaned_row.append('')
                    cleaned_table.append(cleaned_row)
                
                return cleaned_table
            
            return None
            
        except Exception as e:
            logger.warning(f"管道分隔表格解析失败: {e}")
            return None
    
    def _is_markdown_table(self, text: str) -> bool:
        """检查是否为markdown表格格式"""
        lines = text.strip().split('\n')
        if len(lines) < 2:
            return False
        
        # 检查是否有分隔行（包含 --- 的行）
        has_separator = any('---' in line for line in lines)
        # 检查是否有管道符
        has_pipes = any('|' in line for line in lines)
        
        return has_separator and has_pipes
    
    def _parse_markdown_table(self, text: str) -> Optional[List[List[str]]]:
        """解析markdown格式的表格"""
        try:
            lines = [line.strip() for line in text.strip().split('\n') if line.strip()]
            table_data = []
            
            for line in lines:
                if '---' in line:  # 跳过分隔行
                    continue
                
                if '|' in line:
                    # 移除首尾的 |，然后分割
                    cells = [cell.strip() for cell in line.strip('|').split('|')]
                    if cells and any(cell for cell in cells):  # 确保不是空行
                        table_data.append(cells)
            
            return table_data if len(table_data) >= 2 else None
            
        except Exception as e:
            logger.warning(f"Markdown表格解析失败: {e}")
            return None


# 向后兼容的别名
DocumentParser = UnifiedDocumentParser
